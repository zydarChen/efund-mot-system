#!/usr/bin/env python3
"""
生成汇报PPT - 基于易方达PPT模板V5
内容来源: 汇报思路.md + 系统原型
"""
import warnings
warnings.filterwarnings('ignore')

from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
import os

# --- 颜色定义 ---
DARK_BLUE = RGBColor(0, 80, 150)
BRIGHT_BLUE = RGBColor(30, 185, 225)
DARK_GRAY = RGBColor(60, 60, 60)
WHITE = RGBColor(255, 255, 255)
LIGHT_GRAY = RGBColor(128, 128, 128)

BASE_DIR = os.path.dirname(os.path.abspath(__file__))
TEMPLATE_PATH = os.path.join(BASE_DIR, '..', 'data', '易方达ppt模板V5.pptx')
OUTPUT_PATH = os.path.join(BASE_DIR, '..', 'docs', '汇报PPT.pptx')

# --- Layout indices ---
LAYOUT_COVER = 0       # "2_业务介绍" (封面)
LAYOUT_CHAPTER = 9     # "章节页"
LAYOUT_CONTENT = 10    # "文字内容页"
LAYOUT_END = 5         # "结尾页"


def _set_run_style(run, font_name='华文黑体', font_size=14, color=DARK_GRAY, bold=False):
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.color.rgb = color
    run.font.bold = bold
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn('a:ea'))
    if ea is None:
        ea = __import__('lxml.etree', fromlist=['etree']).SubElement(rPr, qn('a:ea'))
    ea.set('typeface', font_name)


def _set_first(tf, text, **kw):
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = kw.get('alignment', PP_ALIGN.LEFT)
    r = p.add_run()
    r.text = text
    _set_run_style(r, kw.get('font_name', '华文黑体'), kw.get('font_size', 14),
                   kw.get('color', DARK_GRAY), kw.get('bold', False))


def _add_para(tf, text, **kw):
    p = tf.add_paragraph()
    p.alignment = kw.get('alignment', PP_ALIGN.LEFT)
    p.level = kw.get('level', 0)
    sb = kw.get('space_before', 0)
    if sb:
        p.space_before = Pt(sb)
    r = p.add_run()
    r.text = text
    _set_run_style(r, kw.get('font_name', '华文黑体'), kw.get('font_size', 14),
                   kw.get('color', DARK_GRAY), kw.get('bold', False))


def _clear_slides(prs):
    """删除所有已有幻灯片"""
    sldIdLst = prs.slides._sldIdLst
    for sldId in list(sldIdLst):
        sldIdLst.remove(sldId)


def slide_cover(prs):
    s = prs.slides.add_slide(prs.slide_layouts[LAYOUT_COVER])
    tf = s.placeholders[1].text_frame
    _set_first(tf, '客户关键服务管理系统', font_size=28, color=WHITE, bold=True)
    _add_para(tf, '（VMOT）', font_size=28, color=WHITE, bold=True)
    _add_para(tf, '', font_size=10, color=WHITE)
    _add_para(tf, '基于客户投资旅程的关键时刻服务管理', font_size=16, color=WHITE)
    tf10 = s.placeholders[10].text_frame
    _set_first(tf10, '易方达基金  |  金融科技部  |  零售客户服务团', font_size=11, color=WHITE)


def slide_chapter(prs, num, title, subtitle=''):
    s = prs.slides.add_slide(prs.slide_layouts[LAYOUT_CHAPTER])
    tf = s.placeholders[0].text_frame
    _set_first(tf, num, font_size=36, color=DARK_BLUE, bold=True)
    _add_para(tf, title, font_size=24, color=DARK_BLUE, bold=True, space_before=6)
    if subtitle:
        tf1 = s.placeholders[1].text_frame
        _set_first(tf1, subtitle, font_size=14, color=DARK_GRAY)


def slide_content(prs, title, items):
    s = prs.slides.add_slide(prs.slide_layouts[LAYOUT_CONTENT])
    ph_ids = [ph.placeholder_format.idx for ph in s.placeholders]
    if 0 in ph_ids:
        _set_first(s.placeholders[0].text_frame, title, font_size=20, color=DARK_BLUE, bold=True)
    tf = s.placeholders[10].text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        kw = {k: v for k, v in item.items() if k != 'text'}
        if i == 0:
            _set_first(tf, item['text'], **kw)
        else:
            _add_para(tf, item['text'], **kw)


def slide_end(prs):
    s = prs.slides.add_slide(prs.slide_layouts[LAYOUT_END])
    _set_first(s.placeholders[10].text_frame, '谢谢', font_size=36, color=DARK_BLUE,
               bold=True, alignment=PP_ALIGN.CENTER)
    ph_ids = [ph.placeholder_format.idx for ph in s.placeholders]
    if 11 in ph_ids:
        tf = s.placeholders[11].text_frame
        _set_first(tf, 'THANK YOU', font_size=14, color=LIGHT_GRAY,
                   alignment=PP_ALIGN.CENTER, font_name='Arial')
        _add_para(tf, '', font_size=8, color=LIGHT_GRAY)
        _add_para(tf, '易方达基金 · 金融科技部 · 零售客户服务团',
                  font_size=11, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)


# ---- 便捷: 生成 item dict ----
def T(text, size=13, color=DARK_GRAY, bold=False, space_before=4, alignment=PP_ALIGN.LEFT):
    return {'text': text, 'font_size': size, 'color': color, 'bold': bold,
            'space_before': space_before, 'alignment': alignment}

def H(text, size=16):
    return T(text, size=size, color=DARK_BLUE, bold=True, space_before=10)

def B(text, size=13):
    return T(text, size=size, color=DARK_BLUE, bold=True, space_before=6)

def P(text, size=12, space_before=4):
    return T(text, size=size, space_before=space_before)

def S(size=4):
    return T('', size=size, space_before=0)

def G(text, size=11):
    return T(text, size=size, color=LIGHT_GRAY, space_before=2)


def main():
    prs = Presentation(TEMPLATE_PATH)
    _clear_slides(prs)

    # ===== 第1页: 封面 =====
    slide_cover(prs)

    # ===== 第2页: 目录 =====
    slide_content(prs, '汇报目录', [
        T('目  录', size=22, color=DARK_BLUE, bold=True, alignment=PP_ALIGN.CENTER, space_before=0),
        S(8),
        T('01    背景与挑战', size=16, color=DARK_BLUE, bold=True, space_before=14),
        T('        金科角色定位 · 业务压力 · 团队发展', size=12, color=LIGHT_GRAY),
        T('02    愿景与方案', size=16, color=DARK_BLUE, bold=True, space_before=14),
        T('        目标愿景 · 业务逻辑 · 系统逻辑', size=12, color=LIGHT_GRAY),
        T('03    系统Demo演示', size=16, color=DARK_BLUE, bold=True, space_before=14),
        T('        MOT总览 · 客户360 · 策略中心 · 策略挖掘 · 触达中心', size=12, color=LIGHT_GRAY),
        T('04    总结与展望', size=16, color=DARK_BLUE, bold=True, space_before=14),
        T('        近期规划 · 未来愿景', size=12, color=LIGHT_GRAY),
    ])

    # ===== 第3页: 章节1 =====
    slide_chapter(prs, '01', '背景与挑战',
                  '金科作为公司大中台板块，受理各板块、部门需求，感受行业与公司的变化')

    # ===== 第4页: 业务板块变革 =====
    slide_content(prs, '业务板块面临新变革要求', [
        H('公司要求谋求新发展、新变革'),
        S(),
        P('►  不仅是AI技术探索，本质是未来业务商业逻辑变化、建立新业务壁垒', size=13, space_before=8),
        S(),
        H('零售业务三大方向同样面临变革'),
        S(),
        P('►  销售服务（总分公司销售服务）—— 主动求变、构建竞争壁垒', size=13, space_before=6),
        P('►  营销服务（线上化产品营销服务）—— 内容生产自动化、精准触达', size=13, space_before=6),
        P('►  客户服务（线上化客户运营服务）—— 服务智能化、体验个性化', size=13, space_before=6),
        S(),
        H('客服产研团队的压力与机遇'),
        S(),
        P('►  过去2年：从寥寥几人发展到10多人，但缺乏质变', size=13, space_before=6),
        P('►  想走得更远、更超前 —— 需要一个能承载团队能力的核心产品', size=13, space_before=6),
    ])

    # ===== 第5页: 过去贡献 =====
    slide_content(prs, '过去工作贡献与思考', [
        H('已有业务基础'),
        S(),
        P('►  投资者之家 —— 整体业务体系开展', size=13, space_before=6),
        P('►  智能客户顾问小易 —— AI客服助手建设', size=13, space_before=6),
        P('►  社区陪伴服务 —— 持续建设，覆盖8个讨论区', size=13, space_before=6),
        P('►  服务触点基础设施 —— 电话、在线、易服务、易陪伴', size=13, space_before=6),
        S(),
        H('核心思考'),
        S(),
        P('►  服务能力分散在不同系统中，缺乏统一管理', size=13, space_before=6),
        P('►  以被动响应为主，缺乏主动服务能力', size=13, space_before=6),
        P('►  客户数据未充分整合，难以做到千人千面', size=13, space_before=6),
        B('►  需要一个顶层设计，将碎片化能力整合为系统化方案'),
    ])

    # ===== 第6页: 章节2 =====
    slide_chapter(prs, '02', '愿景与方案',
                  '从客户找服务到服务找客户，以客户为中心的主动服务体系')

    # ===== 第7页: 目标愿景 =====
    slide_content(prs, '目标愿景：从"客户找服务"到"服务找客户"', [
        H('核心转变'),
        S(),
        T('从被动的客户找服务  ➜  主动的服务找客户', size=15, color=DARK_BLUE, bold=True,
          alignment=PP_ALIGN.CENTER, space_before=10),
        S(),
        P('以客户为中心，基于投资旅程分析客户核心诉求', size=13, space_before=8),
        P('挖掘服务关键时点（Moment of Truth），提供千人千面的服务', size=13, space_before=6),
        S(),
        H('什么是MOT（Moment of Truth）'),
        S(),
        P('►  关键服务时刻：客户投资旅程中对体验产生决定性影响的关键触点', size=13, space_before=6),
        P('►  例如：新客首次购基、大额赎回前、基金净值大跌、持仓产品分红等', size=13, space_before=6),
        P('►  在正确的时间，通过正确的渠道，提供正确的服务内容', size=13, space_before=6),
    ])

    # ===== 第8页: 业务逻辑 =====
    slide_content(prs, '业务逻辑：资管机构客户服务优势', [
        H('客户服务领域触点众多', size=15),
        S(),
        P('►  自建域内：电话、在线、智能客服、易服务、易陪伴'),
        P('►  外部域外：互联网社区、线上运营的内容、活动、工具'),
        S(),
        H('服务内容载体丰富多样', size=15),
        S(),
        P('►  产品营销 · 投资教育 · 宣传策划 · 投研成果'),
        S(),
        H('新生代投资者的变化', size=15),
        S(),
        P('►  Z时代客户线上投资习惯 · 社会高质量发展要求'),
        S(),
        H('结论', size=15),
        T('建立符合未来投资者社会属性变化的、发挥公司已有业务能力优势、', size=12,
          color=DARK_BLUE, space_before=6),
        T('能够构造新业务竞争优势的、非传统人力模式的客户服务体系', size=12,
          color=DARK_BLUE, space_before=2),
    ])

    # ===== 第9页: 系统逻辑 =====
    slide_content(prs, '系统逻辑：AI大脑 + 手脚', [
        H('整体架构：两大核心系统'),
        S(),
        B('AI大脑器官1：客户综合信息管理（CSDP）', size=14),
        S(),
        P('►  整合客户全域数据，构建统一客户视图', size=13, space_before=6),
        P('►  客户标签体系、行为分析、投资画像', size=13),
        P('►  为服务决策提供数据基础', size=13),
        S(),
        B('AI大脑器官2：客户高价值MOT服务管理（VMOT）', size=14),
        S(),
        P('►  基于CSDP的客户数据，识别客户关键服务时刻', size=13, space_before=6),
        P('►  智能匹配最优服务策略，实现服务找客户', size=13),
        P('►  MOT总览 → 客户360 → 策略中心 → 策略挖掘 → 触达中心', size=13),
    ])

    # ===== 第10页: 章节3 =====
    slide_chapter(prs, '03', '系统Demo演示',
                  'VMOT系统原型展示：MOT总览 · 客户360 · 策略中心 · 策略挖掘 · 触达中心')

    # ===== 第11页: MOT总览 =====
    slide_content(prs, 'MOT总览 —— 全局数据驾驶舱', [
        H('功能概述', size=15),
        S(),
        P('MOT总览是系统的核心入口，提供全局视角的数据聚合与可视化展示', size=13, space_before=6),
        S(),
        H('核心模块', size=15),
        S(),
        P('►  顶部指标卡片：活跃MOT数、触达客户数、触达成功率、策略转化率'),
        P('►  MOT触发趋势：7日/30日趋势折线图，实时监控服务触发频率'),
        P('►  渠道分布：各触达渠道（APP推送/短信/邮件/企微/电话）占比分析'),
        P('►  客户生命周期分布：认知/考虑/购买/持有/赎回各阶段客户分布'),
        P('►  策略效果排行：Top策略的执行次数与转化率排名'),
        P('►  最近MOT事件：实时事件流，展示最新触发的MOT服务'),
        P('►  客户满意度趋势：满意度NPS评分变化追踪'),
    ])

    # ===== 第12页: 客户360 =====
    slide_content(prs, '客户360视图 —— 全方位客户洞察', [
        H('设计理念', size=15),
        S(),
        P('整合客户全域数据，构建360度全景画像，覆盖13个维度', size=13, space_before=6),
        S(),
        H('三层信息架构', size=15),
        S(),
        B('第一层：客户概览卡片'),
        P('    头像、姓名、等级、风险偏好、总资产、收益率、持仓天数'),
        S(),
        B('第二层：核心分析看板'),
        P('    客户画像（雷达图）· 投资偏好分析 · 行为轨迹时间线'),
        S(),
        B('第三层：9个Tab详情页'),
        P('    基本信息 · 资产总览 · 持仓明细 · 交易记录 · 行为轨迹'),
        P('    服务记录 · 风险评估 · 标签管理 · MOT时刻'),
    ])

    # ===== 第13页: 策略中心 =====
    slide_content(prs, '策略中心 —— 可视化策略管理', [
        H('功能概述', size=15),
        S(),
        P('提供策略全生命周期管理：创建 → 编排 → 模拟 → 上线 → 监控 → 优化', size=13, space_before=6),
        S(),
        H('核心能力', size=15),
        S(),
        B('►  策略列表管理（CRUD）'),
        P('    支持新增/编辑/删除/启停/搜索/筛选，可视化状态管理'),
        S(),
        B('►  策略画布（Strategy Canvas）'),
        P('    拖拽式流程编排，8种节点类型，SVG贝塞尔曲线连接'),
        P('    支持触发器 → 条件判断 → 渠道分发 → 内容推荐的完整流程'),
        S(),
        B('►  策略模拟'),
        P('    上线前模拟预估覆盖人群、预期效果，降低试错成本'),
    ])

    # ===== 第14页: 策略挖掘 =====
    slide_content(prs, '策略挖掘 —— AI驱动的MOT发现', [
        H('功能概述', size=15),
        S(),
        P('利用AI能力，从多维数据中自动发现客户关键服务时刻（MOT）', size=13, space_before=6),
        S(),
        H('四大挖掘维度', size=15),
        S(),
        P('►  行为分析挖掘：基于客户行为序列，识别异常行为模式'),
        G('    如：频繁查看但未操作、连续登录后沉默等'),
        P('►  生命周期挖掘：基于客户生命周期阶段转换的关键节点'),
        G('    如：新手→成长期、活跃→沉默预警等'),
        P('►  事件驱动挖掘：基于外部事件（市场行情/产品公告等）触发'),
        G('    如：大盘暴跌、基金分红、基金经理变更等'),
        P('►  价值预测挖掘：基于客户价值预测模型，识别潜力客户'),
        G('    如：高净值客户流失预警、潜在升级客户识别等'),
    ])

    # ===== 第15页: MOT推荐引擎 =====
    slide_content(prs, 'MOT推荐引擎 —— 智能策略建议', [
        H('功能概述', size=15),
        S(),
        P('每个策略挖掘页面内置AI推荐面板，实时推荐最优服务策略', size=13, space_before=6),
        S(),
        H('推荐机制', size=15),
        S(),
        P('►  基于当前挖掘维度，智能生成针对性策略建议'),
        P('►  每条推荐包含：策略名称、详细描述、置信度评分、适用客群'),
        P('►  支持"采纳"操作，一键将推荐策略加入策略中心'),
        P('►  支持"忽略"操作，记录反馈以优化后续推荐'),
        S(),
        H('业务价值', size=15),
        S(),
        P('►  降低策略设计门槛，辅助业务人员快速创建高质量策略'),
        P('►  基于数据驱动的推荐，提升策略的精准度与转化率'),
        P('►  形成"挖掘→推荐→采纳→执行→反馈"的闭环优化'),
    ])

    # ===== 第16页: 触达中心 =====
    slide_content(prs, '触达中心 —— 全渠道服务触达', [
        H('功能概述', size=15),
        S(),
        P('统一管理所有客户触达渠道，确保服务在正确的时间通过正确的渠道送达', size=13, space_before=6),
        S(),
        H('四大管理模块', size=15),
        S(),
        B('►  渠道管理'),
        P('    APP推送、短信、邮件、企微、电话等渠道配置与监控'),
        B('►  消息模板'),
        P('    模板CRUD管理，支持变量插入、多渠道适配、版本管理'),
        B('►  触达记录'),
        P('    完整触达日志，支持状态追踪（已发送/已送达/已阅读/已转化）'),
        B('►  效果分析'),
        P('    多维度效果分析，送达率/打开率/转化率对比，A/B测试支持'),
    ])

    # ===== 第17页: 章节4 =====
    slide_chapter(prs, '04', '总结与展望', '近期规划与未来愿景')

    # ===== 第18页: 近期规划 =====
    slide_content(prs, '近期规划与落地路径', [
        H('一期目标（Demo → MVP）', size=15),
        S(),
        P('►  完成系统原型验证，明确核心功能范围与数据需求'),
        P('►  对接CSDP客户数据平台，实现真实客户数据接入'),
        P('►  落地1-2个核心MOT场景（如大额赎回预警、新客引导）'),
        S(),
        H('二期目标（MVP → 生产）', size=15),
        S(),
        P('►  AI策略挖掘引擎上线，实现智能MOT发现与推荐'),
        P('►  全渠道触达打通，与现有客服系统集成'),
        P('►  建立效果评估体系，持续优化策略质量'),
        S(),
        H('远期愿景', size=15),
        S(),
        P('►  构建完整的"CSDP + VMOT"双轮驱动体系'),
        P('►  实现真正的"服务找客户"，提升客户获得感与公司品牌价值'),
    ])

    # ===== 第19页: 结尾 =====
    slide_end(prs)

    os.makedirs(os.path.dirname(OUTPUT_PATH), exist_ok=True)
    prs.save(OUTPUT_PATH)
    print(f'汇报PPT已生成: {OUTPUT_PATH}')
    print(f'共 {len(prs.slides)} 页')


if __name__ == '__main__':
    main()
