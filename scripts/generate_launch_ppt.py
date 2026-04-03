#!/usr/bin/env python3
"""
生成产品发布会PPT - 基于易方达PPT模板V5
面向更广泛受众的产品发布会风格演示
"""
import warnings
warnings.filterwarnings('ignore')

from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
import os

# --- 颜色定义 ---
DARK_BLUE = RGBColor(0, 80, 150)
BRIGHT_BLUE = RGBColor(30, 185, 225)
DARK_GRAY = RGBColor(60, 60, 60)
WHITE = RGBColor(255, 255, 255)
LIGHT_GRAY = RGBColor(128, 128, 128)
ACCENT_ORANGE = RGBColor(230, 140, 30)

BASE_DIR = os.path.dirname(os.path.abspath(__file__))
TEMPLATE_PATH = os.path.join(BASE_DIR, '..', 'data', '易方达ppt模板V5.pptx')
OUTPUT_PATH = os.path.join(BASE_DIR, '..', 'docs', '产品发布会PPT.pptx')

# --- Layout indices ---
LAYOUT_COVER = 0       # "2_业务介绍" (封面)
LAYOUT_CHAPTER = 9     # "章节页"
LAYOUT_CONTENT = 10    # "文字内容页"
LAYOUT_END = 5         # "结尾页"


def _set_run_style(run, font_name='华文黑体', font_size=14, color=DARK_GRAY, bold=False):
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.color.rgb = color
    run.font.bold = bold
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn('a:ea'))
    if ea is None:
        ea = __import__('lxml.etree', fromlist=['etree']).SubElement(rPr, qn('a:ea'))
    ea.set('typeface', font_name)


def _set_first(tf, text, **kw):
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = kw.get('alignment', PP_ALIGN.LEFT)
    r = p.add_run()
    r.text = text
    _set_run_style(r, kw.get('font_name', '华文黑体'), kw.get('font_size', 14),
                   kw.get('color', DARK_GRAY), kw.get('bold', False))


def _add_para(tf, text, **kw):
    p = tf.add_paragraph()
    p.alignment = kw.get('alignment', PP_ALIGN.LEFT)
    p.level = kw.get('level', 0)
    sb = kw.get('space_before', 0)
    if sb:
        p.space_before = Pt(sb)
    r = p.add_run()
    r.text = text
    _set_run_style(r, kw.get('font_name', '华文黑体'), kw.get('font_size', 14),
                   kw.get('color', DARK_GRAY), kw.get('bold', False))


def _clear_slides(prs):
    sldIdLst = prs.slides._sldIdLst
    for sldId in list(sldIdLst):
        sldIdLst.remove(sldId)


def slide_cover(prs):
    s = prs.slides.add_slide(prs.slide_layouts[LAYOUT_COVER])
    tf = s.placeholders[1].text_frame
    _set_first(tf, 'VMOT', font_size=36, color=WHITE, bold=True)
    _add_para(tf, '客户关键服务管理系统', font_size=24, color=WHITE, bold=True, space_before=4)
    _add_para(tf, '', font_size=8, color=WHITE)
    _add_para(tf, '让服务找到客户 · 让每一次触达都有价值', font_size=14, color=WHITE)
    tf10 = s.placeholders[10].text_frame
    _set_first(tf10, '易方达基金 · 产品发布会', font_size=11, color=WHITE)


def slide_chapter(prs, num, title, subtitle=''):
    s = prs.slides.add_slide(prs.slide_layouts[LAYOUT_CHAPTER])
    tf = s.placeholders[0].text_frame
    _set_first(tf, num, font_size=36, color=DARK_BLUE, bold=True)
    _add_para(tf, title, font_size=24, color=DARK_BLUE, bold=True, space_before=6)
    if subtitle:
        tf1 = s.placeholders[1].text_frame
        _set_first(tf1, subtitle, font_size=14, color=DARK_GRAY)


def slide_content(prs, title, items):
    s = prs.slides.add_slide(prs.slide_layouts[LAYOUT_CONTENT])
    ph_ids = [ph.placeholder_format.idx for ph in s.placeholders]
    if 0 in ph_ids:
        _set_first(s.placeholders[0].text_frame, title, font_size=20, color=DARK_BLUE, bold=True)
    tf = s.placeholders[10].text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        kw = {k: v for k, v in item.items() if k != 'text'}
        if i == 0:
            _set_first(tf, item['text'], **kw)
        else:
            _add_para(tf, item['text'], **kw)


def slide_end(prs):
    s = prs.slides.add_slide(prs.slide_layouts[LAYOUT_END])
    _set_first(s.placeholders[10].text_frame, '谢谢', font_size=36, color=DARK_BLUE,
               bold=True, alignment=PP_ALIGN.CENTER)
    ph_ids = [ph.placeholder_format.idx for ph in s.placeholders]
    if 11 in ph_ids:
        tf = s.placeholders[11].text_frame
        _set_first(tf, 'THANK YOU', font_size=14, color=LIGHT_GRAY,
                   alignment=PP_ALIGN.CENTER, font_name='Arial')
        _add_para(tf, '', font_size=8, color=LIGHT_GRAY)
        _add_para(tf, 'VMOT — 让服务找到客户',
                  font_size=11, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)


# ---- 便捷 item 构造 ----
def T(text, size=13, color=DARK_GRAY, bold=False, space_before=4, alignment=PP_ALIGN.LEFT):
    return {'text': text, 'font_size': size, 'color': color, 'bold': bold,
            'space_before': space_before, 'alignment': alignment}

def H(text, size=16):
    return T(text, size=size, color=DARK_BLUE, bold=True, space_before=10)

def B(text, size=13):
    return T(text, size=size, color=DARK_BLUE, bold=True, space_before=6)

def P(text, size=12, space_before=4):
    return T(text, size=size, space_before=space_before)

def S(size=4):
    return T('', size=size, space_before=0)

def G(text, size=11):
    return T(text, size=size, color=LIGHT_GRAY, space_before=2)

def A(text, size=12):
    return T(text, size=size, color=ACCENT_ORANGE, bold=True, space_before=6)


def main():
    prs = Presentation(TEMPLATE_PATH)
    _clear_slides(prs)

    # ===== P1: 封面 =====
    slide_cover(prs)

    # ===== P2: 开场 — 一个故事 =====
    slide_content(prs, '一个基金投资者的故事', [
        H('场景一：恐慌中的沉默', size=15),
        S(),
        P('小王持有的基金净值连续3天下跌超过8%', size=13, space_before=6),
        P('他打开APP看了5次，却没有做任何操作', size=13),
        P('他在等一个声音告诉他"别怕，我们在"', size=13),
        G('—— 但没有人联系他'),
        S(6),
        H('场景二：错过的挽留', size=15),
        S(),
        P('老李是一位持仓3年的高净值客户', size=13, space_before=6),
        P('上周他赎回了60%的持仓，转投了竞品', size=13),
        P('事后团队才发现，他两周前就开始频繁查看赎回页面', size=13),
        G('—— 但没有人注意到这个信号'),
        S(6),
        T('如果我们能在关键时刻，主动找到客户呢？', size=14, color=DARK_BLUE,
          bold=True, alignment=PP_ALIGN.CENTER, space_before=10),
    ])

    # ===== P3: 行业洞察 =====
    slide_content(prs, '行业洞察：客户服务的范式转变', [
        H('从"客户找服务"到"服务找客户"', size=15),
        S(),
        P('传统模式：客户遇到问题 → 拨打电话/发送消息 → 被动响应', size=13, space_before=6),
        P('未来模式：系统识别信号 → 匹配策略 → 主动触达 → 闭环优化', size=13),
        S(6),
        H('三大行业趋势', size=15),
        S(),
        B('趋势一：从标准化服务到千人千面'),
        P('    每个客户的投资旅程不同，服务需求也不同'),
        B('趋势二：从人工判断到数据驱动'),
        P('    用AI替代经验，从海量数据中发现服务机会'),
        B('趋势三：从单渠道到全渠道协同'),
        P('    APP推送、短信、邮件、企微、电话——统一编排，精准触达'),
    ])

    # ===== P4: 产品理念 =====
    slide_chapter(prs, '01', '产品理念',
                  'Moment of Truth — 在关键时刻，让服务主动找到客户')

    # ===== P5: 什么是MOT =====
    slide_content(prs, 'MOT — Moment of Truth 关键服务时刻', [
        H('定义', size=15),
        S(),
        P('MOT是客户投资旅程中，对体验产生决定性影响的关键触点', size=13, space_before=6),
        S(6),
        H('典型MOT场景', size=15),
        S(),
        A('▎ 新客引导'),
        P('    首次购基后的7天引导期，建立信任与认知'),
        A('▎ 大额赎回预警'),
        P('    检测到赎回倾向，主动关怀挽留'),
        A('▎ 市场波动安抚'),
        P('    大盘暴跌时，及时推送专业解读与持仓分析'),
        A('▎ 分红再投资'),
        P('    产品分红时，推荐再投资方案，提升客户粘性'),
        A('▎ 沉默客户唤醒'),
        P('    长期未登录客户，定制化唤醒策略'),
    ])

    # ===== P6: 核心价值 =====
    slide_content(prs, '三大核心价值', [
        H('① 客户获得感提升', size=16),
        S(),
        P('►  在客户最需要的时候出现，而不是在客户不需要时打扰', size=13, space_before=6),
        P('►  提供个性化、有温度的服务体验', size=13),
        S(6),
        H('② 运营效率跃升', size=16),
        S(),
        P('►  AI自动识别MOT时机，替代人工经验判断', size=13, space_before=6),
        P('►  策略画布可视化编排，降低策略配置门槛', size=13),
        S(6),
        H('③ 业务增长驱动', size=16),
        S(),
        P('►  精准触达提升转化率，降低客户流失', size=13, space_before=6),
        P('►  数据驱动的闭环优化，持续提升策略效果', size=13),
    ])

    # ===== P7: 章节 - 产品架构 =====
    slide_chapter(prs, '02', '产品架构',
                  '五大核心模块，覆盖MOT服务全生命周期')

    # ===== P8: 系统全景 =====
    slide_content(prs, 'VMOT 系统全景', [
        T('五大核心模块 · 端到端闭环', size=16, color=DARK_BLUE, bold=True,
          alignment=PP_ALIGN.CENTER, space_before=0),
        S(8),
        B('MOT总览', size=14),
        P('    全局数据驾驶舱，核心指标一目了然', size=12),
        S(),
        B('客户360', size=14),
        P('    全方位客户画像，13维度 · 9大Tab详情 · 雷达图分析', size=12),
        S(),
        B('策略中心', size=14),
        P('    可视化策略编排，拖拽式画布 · CRUD管理 · 上线模拟', size=12),
        S(),
        B('策略挖掘', size=14),
        P('    AI驱动MOT发现，行为分析 · 生命周期 · 事件驱动 · 价值预测', size=12),
        S(),
        B('触达中心', size=14),
        P('    全渠道服务触达，渠道管理 · 消息模板 · 触达记录 · 效果分析', size=12),
    ])

    # ===== P9: MOT总览 =====
    slide_content(prs, '模块一：MOT总览 — 全局数据驾驶舱', [
        H('一屏掌握全局', size=15),
        S(),
        P('►  四大核心指标：活跃MOT数、触达客户数、触达成功率、策略转化率', size=12, space_before=6),
        P('►  MOT触发趋势图：7日/30日维度，实时监控', size=12),
        P('►  渠道分布分析：5大渠道触达占比', size=12),
        P('►  客户生命周期：5阶段客户分布', size=12),
        P('►  策略效果排行：Top策略执行 & 转化排名', size=12),
        P('►  实时事件流：最新MOT触发动态', size=12),
        P('►  满意度趋势：NPS评分追踪', size=12),
        S(6),
        T('设计理念：管理者打开一个页面，就能了解全局服务运营状况',
          size=12, color=DARK_BLUE, bold=True, space_before=8),
    ])

    # ===== P10: 客户360 =====
    slide_content(prs, '模块二：客户360 — 全方位客户洞察', [
        H('三层递进式信息架构', size=15),
        S(6),
        A('▎ 第一层：概览卡片'),
        P('    头像 · 姓名 · 客户等级 · 风险偏好 · 总资产 · 收益率 · 持仓天数'),
        S(),
        A('▎ 第二层：分析看板'),
        P('    客户画像雷达图 · 投资偏好分析 · 行为轨迹时间线'),
        S(),
        A('▎ 第三层：详情Tab页（9个）'),
        P('    基本信息 · 资产总览 · 持仓明细 · 交易记录'),
        P('    行为轨迹 · 服务记录 · 风险评估 · 标签管理 · MOT时刻'),
        S(6),
        T('价值：第一次在一个统一界面看到客户的完整画像',
          size=12, color=DARK_BLUE, bold=True, space_before=6),
    ])

    # ===== P11: 策略中心 =====
    slide_content(prs, '模块三：策略中心 — 可视化策略编排', [
        H('策略全生命周期管理', size=15),
        S(),
        T('创建  →  编排  →  模拟  →  上线  →  监控  →  优化', size=14,
          color=DARK_BLUE, bold=True, alignment=PP_ALIGN.CENTER, space_before=8),
        S(6),
        A('▎ 策略画布（核心亮点）'),
        P('    拖拽式流程编排，所见即所得'),
        P('    8种节点类型：触发器、条件、延时、分流、渠道、内容、标签、结束'),
        P('    SVG贝塞尔曲线连接，节点自动吸附网格'),
        S(),
        A('▎ 策略列表'),
        P('    完整CRUD，支持搜索/筛选/启停/状态管理'),
        S(),
        A('▎ 策略模拟'),
        P('    上线前模拟覆盖人群与预期效果，降低试错成本'),
    ])

    # ===== P12: 策略挖掘 =====
    slide_content(prs, '模块四：策略挖掘 — AI驱动的MOT发现', [
        H('四大挖掘维度', size=15),
        S(),
        A('▎ 行为分析挖掘'),
        P('    识别客户行为序列中的异常模式（频繁查看、沉默预警等）'),
        A('▎ 生命周期挖掘'),
        P('    捕捉客户阶段转换的关键节点（新手→成长、活跃→沉默）'),
        A('▎ 事件驱动挖掘'),
        P('    响应外部事件对客户的影响（大盘暴跌、分红、经理变更）'),
        A('▎ 价值预测挖掘'),
        P('    基于价值模型识别潜力客户（流失预警、升级识别）'),
        S(6),
        H('内置MOT推荐引擎', size=15),
        S(),
        P('►  每个挖掘页面内置AI推荐面板，一键采纳进入策略中心'),
        P('►  形成"挖掘→推荐→采纳→执行→反馈"闭环'),
    ])

    # ===== P13: 触达中心 =====
    slide_content(prs, '模块五：触达中心 — 全渠道精准触达', [
        H('统一渠道管理', size=15),
        S(),
        P('APP推送 · 短信 · 邮件 · 企微 · 电话 — 五大渠道统一编排', size=13, space_before=6),
        S(6),
        A('▎ 渠道管理'),
        P('    各渠道配置、健康监控、优先级策略'),
        A('▎ 消息模板'),
        P('    模板CRUD · 变量插入 · 多渠道适配 · 版本管理'),
        A('▎ 触达记录'),
        P('    全链路日志：已发送 → 已送达 → 已阅读 → 已转化'),
        A('▎ 效果分析'),
        P('    送达率 / 打开率 / 转化率多维对比'),
        P('    A/B测试能力，科学优化触达策略'),
    ])

    # ===== P14: 章节 - 技术架构 =====
    slide_chapter(prs, '03', '技术实现',
                  '现代化前端技术栈，开箱即用的系统原型')

    # ===== P15: 技术栈 =====
    slide_content(prs, '技术架构与栈选型', [
        H('前端技术栈', size=15),
        S(),
        P('►  React 18 + TypeScript 5.6 — 类型安全的现代化框架'),
        P('►  Vite 6 — 极速构建，毫秒级热更新'),
        P('►  Tailwind CSS 3.4 — 原子化样式，设计系统一致性'),
        P('►  Recharts 2.15 — 专业级数据可视化图表'),
        S(6),
        H('核心技术亮点', size=15),
        S(),
        B('策略画布引擎'),
        P('    HTML5 Drag & Drop + SVG Bezier Curves + Grid Snapping'),
        P('    8种节点类型，支持复杂服务流程编排'),
        B('CSS Variables 设计系统'),
        P('    全局设计Token，主题切换，一致的视觉体验'),
        B('组件化架构'),
        P('    高复用组件体系，11,400+行高质量TypeScript代码'),
    ])

    # ===== P16: 设计理念 =====
    slide_content(prs, '产品设计理念', [
        H('以业务人员为中心的设计', size=15),
        S(),
        P('►  零代码策略配置：拖拽画布，业务人员可独立完成策略编排', size=13, space_before=6),
        P('►  可视化数据洞察：图表驱动，复杂数据直观呈现', size=13),
        P('►  AI辅助决策：推荐引擎降低专业门槛', size=13),
        S(6),
        H('设计品质标准', size=15),
        S(),
        P('►  信息层级清晰：3层递进式信息架构', size=13, space_before=6),
        P('►  操作高效：核心操作路径不超过3步', size=13),
        P('►  视觉一致：统一的颜色、字体、间距、组件规范', size=13),
        P('►  响应流畅：微交互动效，0.2-0.3s过渡动画', size=13),
    ])

    # ===== P17: 章节 - 路线图 =====
    slide_chapter(prs, '04', '路线图',
                  '从原型到生产的落地路径')

    # ===== P18: 路线图 =====
    slide_content(prs, '产品路线图', [
        A('▎ 当前阶段：系统原型（已完成）'),
        P('    完成VMOT五大模块的交互原型，验证核心产品逻辑'),
        P('    产出完整的产品需求文档、设计文档、技术方案'),
        S(6),
        A('▎ 一期：MVP验证（Q3 2026）'),
        P('    对接CSDP客户数据平台，接入真实客户数据'),
        P('    落地2-3个核心MOT场景（大额赎回预警、新客引导、市场波动安抚）'),
        P('    完成核心策略画布的生产化改造'),
        S(6),
        A('▎ 二期：生产上线（Q1 2027）'),
        P('    AI策略挖掘引擎上线，全渠道触达系统集成'),
        P('    效果评估闭环建立，支撑10+MOT场景'),
        S(6),
        A('▎ 远期：智能服务生态'),
        P('    CSDP+VMOT双轮驱动，实现真正的"服务找客户"'),
    ])

    # ===== P19: 愿景 =====
    slide_content(prs, '我们的愿景', [
        S(8),
        T('让每一个关键时刻', size=20, color=DARK_BLUE, bold=True,
          alignment=PP_ALIGN.CENTER, space_before=20),
        T('都有温度、有价值', size=20, color=DARK_BLUE, bold=True,
          alignment=PP_ALIGN.CENTER, space_before=6),
        S(12),
        T('在客户最需要的时候出现', size=14, color=DARK_GRAY,
          alignment=PP_ALIGN.CENTER, space_before=10),
        T('用正确的方式，提供正确的服务', size=14, color=DARK_GRAY,
          alignment=PP_ALIGN.CENTER, space_before=6),
        T('让客户感受到"你懂我"', size=14, color=DARK_GRAY,
          alignment=PP_ALIGN.CENTER, space_before=6),
        S(12),
        T('VMOT — 让服务找到客户', size=16, color=BRIGHT_BLUE, bold=True,
          alignment=PP_ALIGN.CENTER, space_before=14),
    ])

    # ===== P20: 结尾 =====
    slide_end(prs)

    os.makedirs(os.path.dirname(OUTPUT_PATH), exist_ok=True)
    prs.save(OUTPUT_PATH)
    print(f'产品发布会PPT已生成: {OUTPUT_PATH}')
    print(f'共 {len(prs.slides)} 页')


if __name__ == '__main__':
    main()
